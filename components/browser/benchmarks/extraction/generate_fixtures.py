#!/usr/bin/env python3
"""Generate the small binary fixtures. The output is byte-for-byte deterministic."""
from __future__ import annotations

import argparse
import io
import zipfile
from datetime import datetime, timezone
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches
ROOT=Path(__file__).parent
FIX=ROOT/'fixtures'

def zip_bytes(files: dict[str, bytes]) -> bytes:
    out=io.BytesIO()
    # Store bytes without compression so zlib versions cannot change fixture hashes.
    with zipfile.ZipFile(out,'w',compression=zipfile.ZIP_STORED) as z:
        for name,data in sorted(files.items()):
            info=zipfile.ZipInfo(name,(1980,1,1,0,0,0)); info.compress_type=zipfile.ZIP_STORED; info.external_attr=0o600<<16
            z.writestr(info,data)
    return out.getvalue()

def pdf(lines: list[str], *, table=False, scanned=False) -> bytes:
    # A one-page PDF with fixed object order, fonts, and metadata-free streams.
    if scanned:
        stream=b"q 120 0 0 40 72 700 cm BI /W 1 /H 1 /CS /RGB /BPC 8 ID \x00\x00\x00 EI Q"
    else:
        ops=['BT /F1 12 Tf 72 740 Td']
        for i,line in enumerate(lines):
            safe=line.replace('\\','\\\\').replace('(','\\(').replace(')','\\)')
            ops.append(('0 -18 Td ' if i else '')+f'({safe}) Tj')
        if table: ops += ['ET','72 650 m 400 650 l 400 610 l 72 610 l h S','BT /F1 10 Tf 80 635 Td (Region   Value) Tj 0 -16 Td (North    42) Tj']
        ops.append('ET'); stream='\n'.join(ops).encode()
    objects=[b'<< /Type /Catalog /Pages 2 0 R >>',b'<< /Type /Pages /Kids [3 0 R] /Count 1 >>',b'<< /Type /Page /Parent 2 0 R /MediaBox [0 0 612 792] /Resources << /Font << /F1 5 0 R >> >> /Contents 4 0 R >>',b'<< /Length %d >>\nstream\n'%len(stream)+stream+b'\nendstream',b'<< /Type /Font /Subtype /Type1 /BaseFont /Helvetica >>']
    out=bytearray(b'%PDF-1.4\n'); offsets=[0]
    for n,obj in enumerate(objects,1): offsets.append(len(out)); out+=f'{n} 0 obj\n'.encode()+obj+b'\nendobj\n'
    xref=len(out); out+=f'xref\n0 {len(objects)+1}\n0000000000 65535 f \n'.encode()
    for off in offsets[1:]: out+=f'{off:010d} 00000 n \n'.encode()
    out+=f'trailer << /Size {len(objects)+1} /Root 1 0 R >>\nstartxref\n{xref}\n%%EOF\n'.encode(); return bytes(out)

def pptx_bytes() -> bytes:
    """Build one valid presentation and normalize its ZIP container."""
    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    box = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(8), Inches(2))
    frame = box.text_frame
    frame.text = "Release Readiness"
    frame.add_paragraph().text = "All bounded checks passed."
    fixed = datetime(2000, 1, 1, tzinfo=timezone.utc)
    presentation.core_properties.created = fixed
    presentation.core_properties.modified = fixed
    raw = io.BytesIO()
    presentation.save(raw)
    with zipfile.ZipFile(io.BytesIO(raw.getvalue())) as source:
        return zip_bytes({name: source.read(name) for name in source.namelist()})


def build() -> dict[str,bytes]:
    docx=zip_bytes({'[Content_Types].xml':b'<Types xmlns="http://schemas.openxmlformats.org/package/2006/content-types"><Default Extension="rels" ContentType="application/vnd.openxmlformats-package.relationships+xml"/><Default Extension="xml" ContentType="application/xml"/><Override PartName="/word/document.xml" ContentType="application/vnd.openxmlformats-officedocument.wordprocessingml.document.main+xml"/></Types>','_rels/.rels':b'<Relationships xmlns="http://schemas.openxmlformats.org/package/2006/relationships"><Relationship Id="r1" Type="http://schemas.openxmlformats.org/officeDocument/2006/relationships/officeDocument" Target="word/document.xml"/></Relationships>','word/document.xml':b'<w:document xmlns:w="http://schemas.openxmlformats.org/wordprocessingml/2006/main"><w:body><w:p><w:r><w:t>Quarterly Operations Brief</w:t></w:r></w:p><w:p><w:r><w:t>North region completed 42 safe deployments.</w:t></w:r></w:p></w:body></w:document>'})
    pptx=pptx_bytes()
    xlsx=zip_bytes({'[Content_Types].xml':b'<Types xmlns="http://schemas.openxmlformats.org/package/2006/content-types"><Default Extension="rels" ContentType="application/vnd.openxmlformats-package.relationships+xml"/><Default Extension="xml" ContentType="application/xml"/><Override PartName="/xl/workbook.xml" ContentType="application/vnd.openxmlformats-officedocument.spreadsheetml.sheet.main+xml"/><Override PartName="/xl/worksheets/sheet1.xml" ContentType="application/vnd.openxmlformats-officedocument.spreadsheetml.worksheet+xml"/></Types>','_rels/.rels':b'<Relationships xmlns="http://schemas.openxmlformats.org/package/2006/relationships"><Relationship Id="r1" Type="http://schemas.openxmlformats.org/officeDocument/2006/relationships/officeDocument" Target="xl/workbook.xml"/></Relationships>','xl/workbook.xml':b'<workbook xmlns="http://schemas.openxmlformats.org/spreadsheetml/2006/main" xmlns:r="http://schemas.openxmlformats.org/officeDocument/2006/relationships"><sheets><sheet name="Metrics" sheetId="1" r:id="r1"/></sheets></workbook>','xl/_rels/workbook.xml.rels':b'<Relationships xmlns="http://schemas.openxmlformats.org/package/2006/relationships"><Relationship Id="r1" Type="http://schemas.openxmlformats.org/officeDocument/2006/relationships/worksheet" Target="worksheets/sheet1.xml"/></Relationships>','xl/worksheets/sheet1.xml':b'<worksheet xmlns="http://schemas.openxmlformats.org/spreadsheetml/2006/main"><sheetData><row r="1"><c r="A1" t="inlineStr"><is><t>Region</t></is></c><c r="B1" t="inlineStr"><is><t>Value</t></is></c></row><row r="2"><c r="A2" t="inlineStr"><is><t>North</t></is></c><c r="B2"><v>42</v></c></row></sheetData></worksheet>'})
    oversized = b"A" * 131_072 + b"\nOVERSIZED_REQUIRED_MARKER\n"
    return {'text.pdf':pdf(['Ordinary PDF Guide','Bounded extraction keeps required content.']),'table.pdf':pdf(['Regional Metrics'],table=True),'scanned.pdf':pdf([],scanned=True),'sample.docx':docx,'sample.pptx':pptx,'sample.xlsx':xlsx,'oversized.txt':oversized}

def main():
    ap=argparse.ArgumentParser(); ap.add_argument('--check',action='store_true'); args=ap.parse_args(); FIX.mkdir(exist_ok=True)
    bad=[]
    for name,data in build().items():
        path=FIX/name
        if args.check:
            if not path.exists() or path.read_bytes()!=data: bad.append(name)
        else: path.write_bytes(data)
    if bad: raise SystemExit('generated fixture drift: '+', '.join(bad))
if __name__=='__main__': main()
